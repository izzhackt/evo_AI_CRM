#!/usr/bin/env python3
"""Extract local text and metadata from downloaded WhatsApp media."""

from __future__ import annotations

import argparse
import hashlib
import json
import os
from pathlib import Path
import subprocess
import tempfile
from typing import Any
import xml.etree.ElementTree as ET
import zipfile

from pypdf import PdfReader
from pptx import Presentation


WORD_TEXT = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}t"


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Локальное извлечение текста из медиа WhatsApp")
    parser.add_argument("--media", type=Path, required=True)
    parser.add_argument("--output", type=Path, required=True)
    return parser.parse_args()


def atomic_text(path: Path, text: str) -> None:
    path.parent.mkdir(parents=True, exist_ok=True, mode=0o700)
    os.chmod(path.parent, 0o700)
    fd, temporary = tempfile.mkstemp(prefix=path.name + ".", suffix=".part", dir=path.parent)
    try:
        os.fchmod(fd, 0o600)
        with os.fdopen(fd, "w", encoding="utf-8") as handle:
            handle.write(text)
            handle.flush()
            os.fsync(handle.fileno())
        os.replace(temporary, path)
    except BaseException:
        try:
            os.unlink(temporary)
        except FileNotFoundError:
            pass
        raise


def extract_pdf(path: Path) -> str:
    reader = PdfReader(str(path))
    pages = []
    for number, page in enumerate(reader.pages, 1):
        pages.append(f"## Страница {number}\n\n{page.extract_text() or ''}")
    return "\n\n".join(pages)


def extract_docx(path: Path) -> str:
    paragraphs: list[str] = []
    with zipfile.ZipFile(path) as archive:
        names = [name for name in archive.namelist() if name == "word/document.xml" or name.startswith("word/header") or name.startswith("word/footer")]
        for name in names:
            root = ET.fromstring(archive.read(name))
            values = [node.text or "" for node in root.iter(WORD_TEXT)]
            if values:
                paragraphs.append(" ".join(values))
    return "\n\n".join(paragraphs)


def extract_pptx(path: Path) -> str:
    presentation = Presentation(str(path))
    slides = []
    for number, slide in enumerate(presentation.slides, 1):
        values = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text:
                values.append(text)
        slides.append(f"## Слайд {number}\n\n" + "\n\n".join(values))
    return "\n\n".join(slides)


def ocr_image(path: Path) -> str:
    with tempfile.TemporaryDirectory(prefix="evo-ocr-") as temporary:
        source = path
        if path.suffix.casefold() in {".heic", ".webp"}:
            converted = Path(temporary) / "converted.png"
            subprocess.run(["sips", "-s", "format", "png", str(path), "--out", str(converted)], check=True, capture_output=True)
            source = converted
        result = subprocess.run(
            ["tesseract", str(source), "stdout", "-l", "rus+eng"],
            check=True,
            capture_output=True,
            text=True,
        )
        return result.stdout


def media_metadata(path: Path) -> str:
    result = subprocess.run(
        ["ffprobe", "-v", "error", "-show_entries", "format=duration,size,format_name", "-of", "json", str(path)],
        check=True,
        capture_output=True,
        text=True,
    )
    return result.stdout


def extractor_for(path: Path):
    suffix = path.suffix.casefold()
    if suffix == ".pdf":
        return "PDF: извлечение текстового слоя", extract_pdf
    if suffix == ".docx":
        return "DOCX: извлечение OOXML", extract_docx
    if suffix == ".pptx":
        return "PPTX: извлечение текста слайдов", extract_pptx
    if suffix in {".jpg", ".jpeg", ".png", ".webp", ".heic"}:
        return "Изображение: локальный OCR Tesseract rus+eng", ocr_image
    if suffix in {".ogg", ".mp4", ".mov"}:
        return "Аудио/видео: метаданные; расшифровка отдельным этапом", media_metadata
    return "Неподдерживаемый формат", lambda _: ""


def main() -> int:
    args = parse_args()
    media = args.media.expanduser().resolve()
    output = args.output.expanduser().resolve()
    output.mkdir(parents=True, exist_ok=True, mode=0o700)
    os.chmod(output, 0o700)
    completed = 0
    skipped = 0
    failed = 0
    formats: dict[str, int] = {}
    for source in sorted((media / "Файлы").glob("*")):
        if not source.is_file():
            continue
        destination = output / f"{source.name}.md"
        source_hash = hashlib.sha256(source.read_bytes()).hexdigest()
        if destination.exists() and f"sha256_источника: {source_hash}" in destination.read_text(encoding="utf-8", errors="replace")[:1000]:
            skipped += 1
            continue
        method, extractor = extractor_for(source)
        formats[source.suffix.casefold()] = formats.get(source.suffix.casefold(), 0) + 1
        try:
            text = extractor(source)
            status = "извлечено" if text.strip() else "нет_текста"
            content = "\n".join([
                "---",
                "тип: извлечённый_текст_медиа",
                f"статус: {status}",
                "доступ_для_клиентского_ии: false",
                f"sha256_источника: {source_hash}",
                "---",
                "",
                f"# Извлечённое содержимое — {source.stem}",
                "",
                f"- Источник: `../Медиа/Файлы/{source.name}`",
                f"- Метод: {method}",
                "",
                "## Содержимое",
                "",
                text.strip() or "Текстовое содержимое не обнаружено.",
                "",
            ])
            atomic_text(destination, content)
            completed += 1
        except Exception as error:
            content = "\n".join([
                "---",
                "тип: извлечённый_текст_медиа",
                "статус: ошибка",
                "доступ_для_клиентского_ии: false",
                f"sha256_источника: {source_hash}",
                "---",
                "",
                f"# Ошибка извлечения — {source.stem}",
                "",
                f"- Источник: `../Медиа/Файлы/{source.name}`",
                f"- Метод: {method}",
                f"- Ошибка: `{type(error).__name__}: {str(error)[:500]}`",
                "",
            ])
            atomic_text(destination, content)
            failed += 1
    print(json.dumps({"извлечено": completed, "пропущено": skipped, "ошибок": failed, "форматы": formats}, ensure_ascii=False))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
